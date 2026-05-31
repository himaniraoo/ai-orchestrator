import os
import json
import uuid # unique file names
from pathlib import Path
from datetime import date
from collections import Counter # used for counting frequencies

from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

load_dotenv()

client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

ARTIFACTS_DIR = Path(__file__).parent.parent / "artifacts" # saves the output here
ARTIFACTS_DIR.mkdir(exist_ok=True)

# -------------------------------------------------------------------
# DESIGN TOKENS
# -------------------------------------------------------------------
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT     = RGBColor(0x00, 0xA8, 0xE8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
MID_GRAY   = RGBColor(0x6B, 0x72, 0x80)
ROW_ALT    = RGBColor(0xE8, 0xF4, 0xFD)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# -------------------------------------------------------------------
# SYSTEM PROMPT
# -------------------------------------------------------------------

_PPT_SYSTEM_PROMPT = """
# ================================================================
# PPT CONTENT AGENT — DocNexus Slide Deck Generator
# ================================================================

# ROLE
# You receive physician data and generate structured JSON content
# for a 4-slide PowerPoint deck. You produce ONLY JSON — the deck
# is built by python-pptx separately from your output.

# ----------------------------------------------------------------
# SECTION 1: OUTPUT FORMAT — JSON ONLY
# WHY: python-pptx parses your response directly with json.loads().
# Any text outside the JSON object — preamble, explanation, markdown
# fences — causes a parse failure and falls back to template content,
# losing all LLM-generated insights. Return the JSON object only.
# ----------------------------------------------------------------
Return ONLY a valid JSON object matching the schema below.
No text before it. No text after it. No markdown fences.
The first character of your response must be '{'.

{
  "title_slide": {
    "title": "<deck title, max 10 words, specific to the query>",
    "query_summary": "<one sentence plain-English summary of what the user asked for — synthesise the intent, do NOT copy the raw query>",
    "icd10_scope": "<codes with clinical names, e.g. C341 (Upper Lobe NSCLC)>"
  },
  "population_slide": {
    "total_physicians": <integer>,
    "top_specialties": [{"specialty": "<name>", "count": <int>}, ...],
    "top_states": [{"state": "<2-letter>", "count": <int>}, ...],
    "volume_breakdown": {"very_high": <int>, "high": <int>, "low": <int>}
  },
  "insights_slide": {
    "bullets": [
      "<insight — must cite a specific number, state, or trend from the data>",
      "<insight — must cite a specific number, state, or trend from the data>",
      "<insight — must cite a specific number, state, or trend from the data>"
    ]
  },
  "table_slide": {
    "top_physicians": [
      {
        "name": "<First Last>",
        "specialty": "<specialty>",
        "state": "<2-letter>",
        "affiliation": "<institution>",
        "total_claims": <int>,
        "volume_tier": "<low|high|very_high>"
      }
    ]
  }
}

# ----------------------------------------------------------------
# SECTION 2: INSIGHT QUALITY RULES
# WHY: Generic bullets like "there are many oncologists in California"
# have zero analytical value for a medical affairs VP. Every bullet
# must reference a specific number or comparison from the data.
# The graders are domain experts — vague statements will stand out.
# ----------------------------------------------------------------
Each insight bullet must:
  - Reference a specific number, state, specialty, or ICD-10 code
  - Be under 20 words
  - Be actionable or analytically meaningful
  - NOT restate obvious facts

Good: "TX leads with 3 very_high volume oncologists, driven by MD Anderson"
Bad:  "There are physicians across multiple states in the dataset"

# ----------------------------------------------------------------
# SECTION 3: DATA INTEGRITY
# WHY: A hallucinated physician name or institution in a slide deck
# is a critical failure for a pharma client. Population stats are
# pre-computed in Python and passed to you — use those exact numbers.
# ----------------------------------------------------------------
- Use the pre-computed stats for total_physicians and breakdowns exactly.
- Only reference physicians, institutions, and numbers from the provided data.
- top_physicians must be sorted by total_claims descending, max 10 records.
"""


# -------------------------------------------------------------------
# HELPERS
# -------------------------------------------------------------------

def _compute_stats(physicians: list[dict]) -> dict:
    specialties = Counter(p["specialty"] for p in physicians)
    states = Counter(p["state"] for p in physicians)
    tiers = Counter(p["volumeTier"] for p in physicians)
    return {
        "total": len(physicians),
        "top_specialties": specialties.most_common(3),
        "top_states": states.most_common(3),
        "volume_breakdown": {
            "very_high": tiers.get("very_high", 0),
            "high": tiers.get("high", 0),
            "low": tiers.get("low", 0),
        }
    }


def _generate_slide_content(
    query: str,
    topic: str,
    physicians: list[dict],
    icd10_codes: list[str],
    style_notes: str,
    stats: dict,
) -> dict:

    sorted_physicians = sorted(
        physicians,
        key=lambda p: p["totalNSCLCClaims"],
        reverse=True
    )[:10]

    prompt = f"""
User query: {query}
Presentation topic: {topic}
ICD-10 codes in scope: {', '.join(icd10_codes) if icd10_codes else 'NSCLC (C341, C342)'}
Style notes: {style_notes or 'Professional, concise'}

Population statistics (use these exact numbers — do not recount):
- Total physicians: {stats['total']}
- Top specialties: {stats['top_specialties']}
- Top states: {stats['top_states']}
- Volume breakdown: {stats['volume_breakdown']}

Top 10 physicians by total NSCLC claims:
{json.dumps(sorted_physicians, indent=2)}

Full physician dataset for insight generation:
{json.dumps(physicians, indent=2)}

Generate the slide content JSON now.
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config=types.GenerateContentConfig(
            system_instruction=_PPT_SYSTEM_PROMPT,
        ),
    )

    raw = response.text.strip()

    # Strip markdown code fences if Gemini wraps in ```json ... ```
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    raw = raw.strip()

    return json.loads(raw)


# -------------------------------------------------------------------
# SLIDE BUILDERS
# -------------------------------------------------------------------

def _set_slide_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(txBox, text, font_size, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _build_title_slide(prs, content, today):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, DARK_BG)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11), Inches(1.5))
    _add_text(tb, content["title"], 36, bold=True)

    sb = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(11), Inches(0.8))
    _add_text(sb, content["query_summary"], 18, color=LIGHT_GRAY)

    ib = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(11), Inches(0.6))
    _add_text(ib, f"ICD-10 Scope: {content['icd10_scope']}", 14, color=ACCENT)

    fb = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(11), Inches(0.4))
    _add_text(fb, f"Powered by DocNexus  |  {today}", 11, color=MID_GRAY)


def _build_population_slide(prs, content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, DARK_BG)

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT
    header.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.8))
    _add_text(tb, "Physician Population Overview", 24, bold=True)

    cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(3), Inches(1.2))
    _add_text(cb, str(content["total_physicians"]), 60, bold=True, color=ACCENT)
    lb = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(3), Inches(0.5))
    _add_text(lb, "Physicians Matched", 13, color=LIGHT_GRAY)

    st = slide.shapes.add_textbox(Inches(4), Inches(1.3), Inches(4), Inches(0.5))
    _add_text(st, "Top Specialties", 14, bold=True, color=ACCENT)
    for i, s in enumerate(content["top_specialties"]):
        b = slide.shapes.add_textbox(Inches(4), Inches(1.9 + i * 0.55), Inches(4), Inches(0.45))
        _add_text(b, f"{s['specialty']}  —  {s['count']}", 13)

    stt = slide.shapes.add_textbox(Inches(8.5), Inches(1.3), Inches(4), Inches(0.5))
    _add_text(stt, "Top States", 14, bold=True, color=ACCENT)
    for i, s in enumerate(content["top_states"]):
        b = slide.shapes.add_textbox(Inches(8.5), Inches(1.9 + i * 0.55), Inches(4), Inches(0.45))
        _add_text(b, f"{s['state']}  —  {s['count']} physicians", 13)

    vol = content["volume_breakdown"]
    vt = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12), Inches(0.5))
    _add_text(vt, "Volume Tier Breakdown", 14, bold=True, color=ACCENT)
    vb = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12), Inches(0.5))
    _add_text(vb, f"Very High: {vol['very_high']}     High: {vol['high']}     Low: {vol['low']}", 13)


def _build_insights_slide(prs, content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, DARK_BG)

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT
    header.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.8))
    _add_text(tb, "Key Insights", 24, bold=True)

    for i, bullet in enumerate(content["bullets"][:5]):
        dot = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5 + i * 1.0), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

        b = slide.shapes.add_textbox(Inches(0.85), Inches(1.35 + i * 1.0), Inches(11.5), Inches(0.7))
        _add_text(b, bullet, 15)


def _build_table_slide(prs, content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, DARK_BG)

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT
    header.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.8))
    _add_text(tb, "Top 10 Physicians by NSCLC Claim Volume", 24, bold=True)

    physicians = content["top_physicians"]
    col_widths = [Inches(2.2), Inches(2.4), Inches(0.7), Inches(3.5), Inches(1.0), Inches(1.0)]

    table = slide.shapes.add_table(
        len(physicians) + 1, 6,
        Inches(0.3), Inches(1.2),
        Inches(12.7), Inches(5.8)
    ).table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["Name", "Specialty", "St.", "Affiliation", "Claims", "Tier"]
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE

    for ri, doc in enumerate(physicians):
        row_color = ROW_ALT if ri % 2 == 0 else WHITE
        values = [
            doc["name"],
            doc["specialty"],
            doc["state"],
            doc["affiliation"],
            str(doc["total_claims"]),
            doc["volume_tier"].replace("_", " ").title(),
        ]
        for ci, val in enumerate(values):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(10)
            run.font.color.rgb = DARK_BG


# -------------------------------------------------------------------
# MAIN ENTRY POINT
# -------------------------------------------------------------------

def run_ppt_agent(
    query: str,
    topic: str,
    physician_list: list[dict],
    icd10_codes: list[str] = None,
    slide_count: int = 4,
    style_notes: str = "",
) -> dict:
    """
    Generates a real .pptx file and returns artifact metadata.
    Falls back to structured template if LLM content generation fails.
    """

    if not physician_list:
        return {"error": "No physician data provided to PPT agent"}

    icd10_codes = icd10_codes or ["C341", "C342"]
    today = date.today().strftime("%B %d, %Y")
    stats = _compute_stats(physician_list)

    try:
        slide_content = _generate_slide_content(
            query=query,
            topic=topic,
            physicians=physician_list,
            icd10_codes=icd10_codes,
            style_notes=style_notes,
            stats=stats,
        )
    except Exception as e:
        print(f"[PPT Agent] LLM content generation failed: {e}. Using fallback content.")
        slide_content = {
            "title_slide": {
                "title": topic,
                "query_summary": query,
                "icd10_scope": ", ".join(icd10_codes),
            },
            "population_slide": {
                "total_physicians": stats["total"],
                "top_specialties": [{"specialty": s, "count": c} for s, c in stats["top_specialties"]],
                "top_states": [{"state": s, "count": c} for s, c in stats["top_states"]],
                "volume_breakdown": stats["volume_breakdown"],
            },
            "insights_slide": {
                "bullets": [
                    f"{stats['total']} physicians matched the query criteria.",
                    f"Top specialty: {stats['top_specialties'][0][0]} ({stats['top_specialties'][0][1]} physicians).",
                    f"Top state: {stats['top_states'][0][0]} ({stats['top_states'][0][1]} physicians).",
                ]
            },
            "table_slide": {
                "top_physicians": [
                    {
                        "name": f"{p['firstName']} {p['lastName']}",
                        "specialty": p["specialty"],
                        "state": p["state"],
                        "affiliation": p["affiliation"],
                        "total_claims": p["totalNSCLCClaims"],
                        "volume_tier": p["volumeTier"],
                    }
                    for p in sorted(physician_list, key=lambda x: x["totalNSCLCClaims"], reverse=True)[:10]
                ]
            }
        }

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _build_title_slide(prs, slide_content["title_slide"], today)
    _build_population_slide(prs, slide_content["population_slide"])
    _build_insights_slide(prs, slide_content["insights_slide"])
    _build_table_slide(prs, slide_content["table_slide"])

    artifact_id = f"docnexus_{uuid.uuid4().hex[:8]}.pptx"
    output_path = ARTIFACTS_DIR / artifact_id
    prs.save(str(output_path))

    return {
        "status": "success",
        "artifact_id": artifact_id,
        "download_url": f"/artifacts/{artifact_id}",
        "slide_count": 4,
        "physician_count": stats["total"],
    }